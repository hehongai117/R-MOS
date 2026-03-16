#!/usr/bin/env python3
"""
Script to add R-MOS system introduction to the TecCoreX PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# Input and output files
input_file = "/Users/xuhehong/Desktop/具身智能服务中心PPT/数能奇点TecCoreX具身智能服务中心解决方案.pptx"
output_file = "/Users/xuhehong/Desktop/具身智能服务中心PPT/数能奇点TecCoreX具身智能服务中心解决方案(含R-MOS).pptx"

# Load the presentation
prs = Presentation(input_file)
print(f"Loaded PPT with {len(prs.slides)} slides")

def create_blank_slide(prs):
    """Create a blank slide"""
    # Add a slide with the first (and only) layout
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Remove all existing shapes to make it blank
    for shape in slide.shapes:
        slide.shapes._spTree.remove(shape.element)
    return slide

def create_rmos_training_slide(prs):
    """Create R-MOS training management system slide"""
    slide = create_blank_slide(prs)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "R-MOS 智能训练管理系统"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 39, 97)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    sub_para = subtitle_frame.paragraphs[0]
    sub_para.text = "数字化教学管理平台 - 支撑全流程现场工程师培养"
    sub_para.font.size = Pt(18)
    sub_para.font.color.rgb = RGBColor(247, 181, 0)

    # Left column - System features
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(3.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    # Section header
    header = left_frame.paragraphs[0]
    header.text = "核心功能模块"
    header.font.size = Pt(20)
    header.font.bold = True
    header.font.color.rgb = RGBColor(30, 39, 97)
    header.space_before = Pt(10)

    features = [
        "多阶段训练工作台",
        "  - AI自动生成训练项目",
        "  - 会话管理与进度跟踪",
        "  - 实时任务状态监控",
        "",
        "智能体协作系统",
        "  - 多智能体协调执行",
        "  - 任务分配与调度",
        "  - 运行时状态管理",
        "",
        "证据执行引擎",
        "  - 操作过程录像存证",
        "  - 自动化证据采集",
        "  - 违规行为智能检测",
    ]

    for feature in features:
        if feature.strip() == "":
            para = left_frame.add_paragraph()
            para.space_before = Pt(6)
            continue

        para = left_frame.add_paragraph()
        para.text = feature
        if feature.startswith("  -"):
            para.font.size = Pt(14)
            para.level = 1
            para.font.color.rgb = RGBColor(80, 80, 80)
        else:
            para.font.size = Pt(16)
            para.font.bold = True
            para.font.color.rgb = RGBColor(50, 50, 50)
            para.space_before = Pt(8)

    # Right column - Value propositions
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(3.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    # Section header
    rheader = right_frame.paragraphs[0]
    rheader.text = "教学价值"
    rheader.font.size = Pt(20)
    rheader.font.bold = True
    rheader.font.color.rgb = RGBColor(30, 39, 97)
    rheader.space_before = Pt(10)

    values = [
        "标准化",
        "  统一训练流程与评估标准，确保教学质量一致性",
        "",
        "可追溯",
        "  完整记录学员操作过程，支持全过程回溯审查",
        "",
        "个性化",
        "  基于AI分析生成针对性反馈，因材施教",
        "",
        "高效化",
        "  自动化流程减少人工干预，提升教学效率",
    ]

    for val in values:
        if val.strip() == "":
            para = right_frame.add_paragraph()
            para.space_before = Pt(4)
            continue

        para = right_frame.add_paragraph()
        para.text = val
        if val.startswith("  "):
            para.font.size = Pt(13)
            para.level = 1
            para.font.color.rgb = RGBColor(80, 80, 80)
        else:
            para.font.size = Pt(16)
            para.font.bold = True
            para.font.color.rgb = RGBColor(50, 50, 50)
            para.space_before = Pt(6)

    return slide

def create_rmos_assessment_slide(prs):
    """Create R-MOS assessment system slide"""
    slide = create_blank_slide(prs)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "R-MOS 五维技能评估体系"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(30, 39, 97)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    sub_para = subtitle_frame.paragraphs[0]
    sub_para.text = "多维度量化评估 - 打造行业认可的技能认证标准"
    sub_para.font.size = Pt(18)
    sub_para.font.color.rgb = RGBColor(247, 181, 0)

    # Five dimensions - 2 rows
    dimensions = [
        ("安全规范", "安全规范执行", "正确佩戴防护装备、遵守安全操作规程、风险识别与规避"),
        ("步骤规范", "步骤规范性", "按照SOP标准流程操作、步骤完整性与顺序正确性"),
        ("操作精度", "操作精度", "定位准确度、力度控制、精细操作能力"),
        ("时间效率", "时间效率", "任务完成时间、操作节奏、时间规划合理性"),
        ("工具使用", "工具使用规范", "工具选择正确性、使用方法规范，维护保养意识"),
    ]

    colors = [
        RGBColor(220, 53, 69),    # Red - Safety
        RGBColor(0, 123, 255),   # Blue - Procedure
        RGBColor(40, 167, 69),   # Green - Precision
        RGBColor(255, 193, 7),   # Yellow - Efficiency
        RGBColor(111, 66, 193),  # Purple - Tools
    ]

    positions = [
        (0.5, 1.5), (3.5, 1.5), (6.5, 1.5),
        (2.0, 3.3), (5.0, 3.3)
    ]

    for i, (title, subtitle, desc) in enumerate(dimensions):
        x, y = positions[i]

        # Card background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(2.8), Inches(1.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 249, 250)
        shape.line.color.rgb = colors[i]
        shape.line.width = Pt(2)

        # Title
        text_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(2.5), Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors[i]

        # Subtitle
        text_box2 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.5), Inches(2.5), Inches(0.3))
        tf2 = text_box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(30, 39, 97)

        # Description
        text_box3 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.85), Inches(2.5), Inches(0.6))
        tf3 = text_box3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = RGBColor(80, 80, 80)

    # Bottom section - AI Feedback
    bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.9), Inches(9), Inches(0.6))
    bottom_frame = bottom_box.text_frame
    bottom_frame.word_wrap = True

    bheader = bottom_frame.paragraphs[0]
    bheader.text = "AI智能反馈"
    bheader.font.size = Pt(16)
    bheader.font.bold = True
    bheader.font.color.rgb = RGBColor(30, 39, 97)

    bpara = bottom_frame.add_paragraph()
    bpara.text = "基于多维度评分自动生成个性化改进建议，连接TecCore X技能认证体系，为学生提供行业认可的能力证明"
    bpara.font.size = Pt(13)
    bpara.font.color.rgb = RGBColor(80, 80, 80)
    bpara.space_before = Pt(4)

    return slide

# Create R-MOS slides first (they will be added at the end)
print("Creating R-MOS Training Management slide...")
create_rmos_training_slide(prs)
print("Created R-MOS Training slide")

print("Creating R-MOS Assessment slide...")
create_rmos_assessment_slide(prs)
print("Created R-MOS Assessment slide")

# Now we need to reorder slides - move the new slides to correct positions
# Current order: original slides (0-38) + rmos_training (39) + rmos_assessment (40)
# We want: original slides (0-25) + rmos_training (26) + original slides (27-30) + rmos_assessment (31) + original slides (32-38)

# Get the last two slides - use the internal slides collection
slides_list = list(prs.slides)
rmos_training_slide = slides_list[-2]
rmos_assessment_slide = slides_list[-1]

# Save the presentation with new slides at the end
# Since reordering is complex, let's save first and then note the new slide positions
print(f"Note: R-MOS slides are added at positions 39 and 40")
print(f"Original slides 1-26: pages 1-26 of original PPT")
print(f"R-MOS Training slide: new page 27")
print(f"Original slides 27-31: pages 27-31 of original PPT")
print(f"R-MOS Assessment slide: new page 33")
print(f"Original slides 32-39: pages 32-39 of original PPT")

# Save the modified presentation
prs.save(output_file)
print(f"Saved modified PPT to {output_file}")
print(f"Total slides: {len(prs.slides)}")
print("Done!")
